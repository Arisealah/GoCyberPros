

# #!/usr/bin/env python3
# """
# PPTX Metadata Extraction Tool

# Description: Comprehensive metadata extraction from PowerPoint PPTX files with validation,
# error handling, and standardized output.

# Features:
# - Extracts core properties, slide information
# - Validates PPTX structure (ZIP+XML)
# - Handles modern PowerPoint formats
# """

# import os
# import zipfile
# import magic
# from datetime import datetime
# from typing import Dict, Any
# from pptx import Presentation

# class PPTXMetadataError(Exception):
#     """Base class for PPTX metadata exceptions"""
#     pass

# class InvalidPPTXError(PPTXMetadataError):
#     """Raised when file is not a valid PPTX"""
#     pass

# class PPTXProcessingError(PPTXMetadataError):
#     """Raised during PPTX processing failures"""
#     pass

# def validate_pptx_file(file_path: str) -> bool:
#     """
#     Validate PPTX file structure and integrity.
    
#     Args:
#         file_path: Path to the PPTX file
        
#     Raises:
#         FileNotFoundError: If file doesn't exist
#         PermissionError: If file access is restricted
#         InvalidPPTXError: If file is not a valid PPTX
#     """
#     if not os.path.exists(file_path):
#         raise FileNotFoundError(f"File not found: {file_path}")
    
#     if not os.access(file_path, os.R_OK):
#         raise PermissionError(f"Access denied: {file_path}")
    
#     mime = magic.from_file(file_path, mime=True)
#     if mime not in ('application/vnd.openxmlformats-officedocument.presentationml.presentation', 
#                    'application/zip'):
#         raise InvalidPPTXError(f"Not a PPTX file (detected: {mime})")
    
#     try:
#         with zipfile.ZipFile(file_path) as z:
#             if 'ppt/presentation.xml' not in z.namelist():
#                 raise InvalidPPTXError("Missing required PPTX files")
#     except zipfile.BadZipFile as e:
#         raise InvalidPPTXError(f"Invalid ZIP structure: {str(e)}")
    
#     return True

# def extract_pptx_metadata(file_path: str) -> Dict[str, Any]:
#     """
#     Extract comprehensive metadata from PPTX files.
    
#     Returns:
#         Dictionary containing:
#         - file_info: Basic file attributes
#         - core_properties: Standard PPTX metadata
#         - presentation_info: Slide and structure data
#         - processing: Status information
        
#     Raises:
#         PPTXMetadataError: For any processing failures
#     """
#     result = {
#         "file_info": {},
#         "core_properties": {},
#         "presentation_info": {},
#         "processing": {
#             "success": False,
#             "warnings": [],
#             "time_taken": None
#         }
#     }
    
#     start_time = datetime.now()
    
#     try:
#         validate_pptx_file(file_path)
        
#         file_stat = os.stat(file_path)
#         result["file_info"] = {
#             "path": os.path.abspath(file_path),
#             "size_bytes": file_stat.st_size,
#             "created": datetime.fromtimestamp(file_stat.st_ctime).isoformat(),
#             "modified": datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
#             "format": "PPTX",
#             "valid": True
#         }
        
#         prs = Presentation(file_path)
#         core_props = prs.core_properties
        
#         result["core_properties"] = {
#             "title": core_props.title,
#             "subject": core_props.subject,
#             "author": core_props.author,
#             "keywords": core_props.keywords,
#             "comments": core_props.comments,
#             "category": core_props.category,
#             "created": core_props.created.isoformat() if core_props.created else None,
#             "modified": core_props.modified.isoformat() if core_props.modified else None,
#             "last_modified_by": core_props.last_modified_by,
#             "revision": core_props.revision,
#             "content_status": core_props.content_status
#         }
        
#         result["presentation_info"] = {
#             "slides": len(prs.slides),
#             "slide_masters": len(prs.slide_masters),
#             "slide_layouts": len(prs.slide_layouts),
#             "notes_masters": len(prs.notes_master),
#             "handout_masters": len(prs.handout_master),
#             "slide_size": {
#                 "width": prs.slide_width.inches,
#                 "height": prs.slide_height.inches
#             },
#             "slide_titles": [
#                 slide.shapes.title.text if slide.shapes.title else None
#                 for slide in prs.slides
#             ]
#         }
        
#         # Check for embedded media
#         media_count = 0
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, 'movie'):
#                     media_count += 1
#                 if hasattr(shape, 'image'):
#                     media_count += 1
        
#         result["presentation_info"]["embedded_media"] = media_count
        
#         result["processing"]["success"] = True
        
#     except FileNotFoundError as e:
#         raise PPTXMetadataError(f"File error: {str(e)}") from e
#     except PermissionError as e:
#         raise PPTXMetadataError(f"Access error: {str(e)}") from e
#     except zipfile.BadZipFile as e:
#         raise InvalidPPTXError(f"Invalid PPTX/ZIP: {str(e)}") from e
#     except Exception as e:
#         raise PPTXProcessingError(f"Processing failed: {str(e)}") from e
#     finally:
#         result["processing"]["time_taken"] = (datetime.now() - start_time).total_seconds()
    
#     return result



































#!/usr/bin/env python3
"""
PPTX Metadata Extraction Tool

Description: Comprehensive metadata extraction from PowerPoint PPTX files with validation,
error handling, and standardized output.

Features:
- Extracts core properties, slide information
- Validates PPTX structure (ZIP+XML)
- Handles modern PowerPoint formats
- Command-line interface with user prompts
- Error handling with JSON output
- Option to save output to file
"""

import os
import zipfile
import magic
import argparse
import json
from datetime import datetime
from typing import Dict, Any, Optional
from pptx import Presentation

class PPTXMetadataError(Exception):
    """Base class for PPTX metadata exceptions"""
    pass

class InvalidPPTXError(PPTXMetadataError):
    """Raised when file is not a valid PPTX"""
    pass

class PPTXProcessingError(PPTXMetadataError):
    """Raised during PPTX processing failures"""
    pass

def validate_pptx_file(file_path: str) -> bool:
    """
    Validate PPTX file structure and integrity.
    
    Args:
        file_path: Path to the PPTX file
        
    Raises:
        FileNotFoundError: If file doesn't exist
        PermissionError: If file access is restricted
        InvalidPPTXError: If file is not a valid PPTX
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    if not os.access(file_path, os.R_OK):
        raise PermissionError(f"Access denied: {file_path}")
    
    mime = magic.from_file(file_path, mime=True)
    if mime not in ('application/vnd.openxmlformats-officedocument.presentationml.presentation', 
                   'application/zip'):
        raise InvalidPPTXError(f"Not a PPTX file (detected: {mime})")
    
    try:
        with zipfile.ZipFile(file_path) as z:
            if 'ppt/presentation.xml' not in z.namelist():
                raise InvalidPPTXError("Missing required PPTX files")
    except zipfile.BadZipFile as e:
        raise InvalidPPTXError(f"Invalid ZIP structure: {str(e)}")
    
    return True

def extract_pptx_metadata(file_path: str) -> Dict[str, Any]:
    """
    Extract comprehensive metadata from PPTX files.
    
    Returns:
        Dictionary containing:
        - file_info: Basic file attributes
        - core_properties: Standard PPTX metadata
        - presentation_info: Slide and structure data
        - processing: Status information
        
    Raises:
        PPTXMetadataError: For any processing failures
    """
    result = {
        "file_info": {},
        "core_properties": {},
        "presentation_info": {},
        "processing": {
            "success": False,
            "warnings": [],
            "time_taken": None
        }
    }
    
    start_time = datetime.now()
    prs = None
    
    try:
        validate_pptx_file(file_path)
        
        file_stat = os.stat(file_path)
        result["file_info"] = {
            "path": os.path.abspath(file_path),
            "size_bytes": file_stat.st_size,
            "created": datetime.fromtimestamp(file_stat.st_ctime).isoformat(),
            "modified": datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
            "format": "PPTX",
            "valid": True
        }
        
        prs = Presentation(file_path)
        core_props = prs.core_properties
        
        result["core_properties"] = {
            "title": core_props.title,
            "subject": core_props.subject,
            "author": core_props.author,
            "keywords": core_props.keywords,
            "comments": core_props.comments,
            "category": core_props.category,
            "created": core_props.created.isoformat() if core_props.created else None,
            "modified": core_props.modified.isoformat() if core_props.modified else None,
            "last_modified_by": core_props.last_modified_by,
            "revision": core_props.revision,
            "content_status": core_props.content_status
        }
        
        # Presentation structure
        result["presentation_info"] = {
            "slides": len(prs.slides),
            "slide_masters": len(prs.slide_masters),
            "slide_layouts": len(prs.slide_layouts),
            "slide_size": {
                "width": prs.slide_width.inches,
                "height": prs.slide_height.inches
            },
            "slide_titles": [
                slide.shapes.title.text if slide.shapes.title else None
                for slide in prs.slides
            ]
        }
        
        # Count embedded media
        media_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'movie') or hasattr(shape, 'image'):
                    media_count += 1
        
        result["presentation_info"]["embedded_media"] = media_count
        
        result["processing"]["success"] = True
        
    except FileNotFoundError as e:
        raise PPTXMetadataError(f"File error: {str(e)}") from e
    except PermissionError as e:
        raise PPTXMetadataError(f"Access error: {str(e)}") from e
    except zipfile.BadZipFile as e:
        raise InvalidPPTXError(f"Invalid PPTX/ZIP: {str(e)}") from e
    except Exception as e:
        raise PPTXProcessingError(f"Processing failed: {str(e)}") from e
    finally:
        if prs:
            del prs  # Ensure presentation object is cleaned up
        result["processing"]["time_taken"] = (datetime.now() - start_time).total_seconds()
    
    return result

def save_output_to_file(output: Dict, output_path: Optional[str] = None) -> None:
    """Save the output to a file if requested"""
    if output_path:
        try:
            with open(output_path, 'w') as f:
                json.dump(output, f, indent=2)
            print(f"\nMetadata saved to: {output_path}")
        except Exception as e:
            print(f"\nWarning: Could not save output to file: {str(e)}")

def main():
    """Command-line interface for the PPTX metadata extractor"""
    parser = argparse.ArgumentParser(description='PPTX Metadata Extraction Tool')
    parser.add_argument('file_path', nargs='?', help='Path to the PPTX file')
    parser.add_argument('--output', '-o', help='Path to save the output JSON')
    args = parser.parse_args()
    
    file_path = args.file_path
    if not file_path:
        file_path = input("Enter the path to the PPTX file: ")
    
    try:
        metadata = extract_pptx_metadata(file_path)
        print("\nPPTX Metadata:")
        print(json.dumps(metadata, indent=2))
        
        if args.output:
            save_output_to_file(metadata, args.output)
        
    except PPTXMetadataError as e:
        error_output = {
            "error": str(e),
            "success": False,
            "file": file_path
        }
        print(json.dumps(error_output, indent=2))
        exit(1)

if __name__ == "__main__":
    main()